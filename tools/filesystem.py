"""Filesystem tools — IO Layer for file operations.

Path resolution is delegated to Workspace (Discovery Layer).
Filesystem only does the actual read/write operations.

Changes:
- _resolve_path now delegates to Workspace first
- read_file/write_file/list_directory record context via WorkspaceService
- Fully backward compatible (fallback to existing logic)
"""
from pathlib import Path
import shutil
from typing import Optional

from agent.registry.tool_registry import registry
from agent.security import (
    is_internal_storage_path,
    is_sensitive_path,
    redact_sensitive_text,
)

# ── Constants ──

ROOT = Path(__file__).parent.parent.resolve()
_working_directory: str = "."  # relative to ROOT
_path_cache: dict = {}

# Cache for the workspace service instance
_workspace_service = None


def _ensure_workspace_path(path: Path, requested: str = "") -> Path:
    """所有文件工具统一限制在当前 workspace 根目录内。"""
    workspace_root = ROOT.resolve()
    resolved = path.resolve()
    try:
        resolved.relative_to(workspace_root)
    except ValueError as exc:
        label = requested or str(path)
        raise PermissionError(f"路径超出 workspace 范围: {label}") from exc
    return resolved


def _get_workspace_service():
    """Lazy-load workspace service."""
    global _workspace_service
    if _workspace_service is None:
        try:
            from agent.services.workspace_service import get_workspace_service
            _workspace_service = get_workspace_service()
        except Exception:
            _workspace_service = None
    return _workspace_service


# ── Working Directory ──


def get_working_directory() -> str:
    """Get current working directory (relative to project root)."""
    return _working_directory


def set_working_directory(path: str) -> str:
    """Set the current working directory for relative path resolution.

    Args:
        path: relative path from project root (e.g. "output", "src")

    Returns:
        Absolute path of the new working directory
    """
    global _working_directory
    target = str(path).strip().rstrip("/")  # 防御：工具参数可能被传 PosixPath
    full = _ensure_workspace_path(ROOT / target, target)

    if not full.exists() or not full.is_dir():
        return f"错误：目录不存在 {path}（项目根目录: {ROOT}）"

    try:
        _working_directory = str(full.relative_to(ROOT))
    except ValueError:
        _working_directory = str(full)

    _path_cache.clear()
    return f"当前工作目录: {_working_directory}"


# ── Path Resolution ──


def _resolve_path(path: str) -> Path:
    """Resolve a path spec to an actual filesystem path.

    Resolution order:
    1. Try Workspace.exact() resolution (uses indexed file tree)
    2. Relative to current working directory
    3. Relative to project root
    4. Common prefixes (input/, src/, output/, docs/)
    5. Fuzzy match via Workspace (best match if unique)
    6. Recursive search fallback
    """
    path = str(path).strip()

    # Check cache first
    if path in _path_cache:
        cached = _path_cache[path]
        if cached.exists():
            return cached

    # Pathological inputs: root references must map to the project root.
    # Workspace fuzzy matcher treats "." as a filename prefix (e.g. ".gitignore"),
    # which would resolve list_directory(".") to a random file instead of the root.
    if path in ("", ".", "./"):
        return ROOT

    # Absolute paths produced by Workspace resolution are already canonical.
    # Check them before fuzzy resolution, otherwise a directory such as
    # ``/project/tests`` may be mistaken for ``tests/__init__.py``.
    absolute = Path(path)
    if absolute.is_absolute():
        return absolute.resolve()

    # Strategy 1: Workspace resolution (exact + fuzzy)
    ws = _get_workspace_service()
    if ws is not None:
        try:
            matches = ws.resolve(path)
        except Exception:
            # 工具可在 bootstrap 前被单独调用；回退到确定性路径解析。
            matches = []
        if matches:
            # Use highest-scored match
            best = matches[0]
            if best.path.exists():
                _path_cache[path] = best.path
                return best.path

    # Strategy 2: Relative to working directory
    if _working_directory != ".":
        cwd_full = (ROOT / _working_directory / path).resolve()
        if cwd_full.exists():
            _path_cache[path] = cwd_full
            return cwd_full

    # Strategy 3: Relative to project root
    full = (ROOT / path).resolve()
    if full.exists():
        _path_cache[path] = full
        return full

    # Strategy 4: Common prefixes
    alternatives = [
        ROOT / "input" / path,
        ROOT / "src" / path,
        ROOT / "output" / path,
        ROOT / "data" / path,
        ROOT / "docs" / path,
        ROOT / path,
    ]
    if _working_directory != ".":
        cwd = ROOT / _working_directory
        alternatives = [
            cwd / "input" / path,
            cwd / "src" / path,
            cwd / "output" / path,
            cwd / "data" / path,
            cwd / path,
        ] + alternatives

    for alt in alternatives:
        if alt.exists():
            _path_cache[path] = alt
            return alt

    # Strategy 5: Recursive search (max 2 levels deep)
    path_name = Path(path).name
    search_roots = [ROOT]
    if _working_directory != ".":
        search_roots.insert(0, ROOT / _working_directory)

    for search_root in search_roots:
        if path_name != path:
            for p in search_root.rglob(path_name):
                if p.is_file():
                    _path_cache[path] = p
                    return p
        else:
            for level1 in search_root.iterdir():
                if level1.name.startswith("."):
                    continue
                if level1.is_dir():
                    target = level1 / path_name
                    if target.exists():
                        _path_cache[path] = target
                        return target
                    for level2 in level1.iterdir():
                        if level2.name.startswith("."):
                            continue
                        if level2.is_dir():
                            target = level2 / path_name
                            if target.exists():
                                _path_cache[path] = target
                                return target

    return full  # Return original path as fallback


def _resolve_operation_path(path: str, *, exact: bool = False) -> Path:
    """Resolve a mutation target without fuzzy redirection when requested.

    Reads may use Workspace discovery.  Mutations must be able to prove which
    path they touched, so File Ops callers pass ``exact=True`` and resolve
    relative to the workspace root directly.
    """
    requested = Path(str(path).strip())
    candidate = requested if requested.is_absolute() else ROOT / requested
    return _ensure_workspace_path(
        candidate if exact else _resolve_path(str(path)),
        str(path),
    )


def _protect_operation_path(path: str, *, label: str) -> Optional[str]:
    """Return a stable error for paths that are not user workspace files."""
    if is_internal_storage_path(path):
        return f"错误：PROTECTED_INTERNAL_PATH：禁止对内部存储执行 {label} 操作。"
    if is_sensitive_path(path):
        return f"错误：出于安全原因，禁止对敏感文件执行 {label} 操作。"
    return None


def _record_edit(path: Path) -> None:
    ws = _get_workspace_service()
    if ws is not None:
        try:
            ws.record_edit(str(path.relative_to(ROOT)))
        except ValueError:
            pass


# ── Tool Implementations ──


def read_file(path: str) -> str:
    """Read file content with automatic path resolution.

    Args:
        path: file path (relative to current directory or project root)

    Returns:
        File content as string
    """
    if is_internal_storage_path(path):
        return "错误：PROTECTED_INTERNAL_PATH：Agent 内部存储只能通过专用接口访问。"
    if is_sensitive_path(path):
        return "错误：出于安全原因，禁止读取敏感文件。"

    try:
        full = _ensure_workspace_path(_resolve_path(path), str(path))
    except PermissionError as e:
        return f"错误：{e}"

    if is_internal_storage_path(full):
        return "错误：PROTECTED_INTERNAL_PATH：Agent 内部存储只能通过专用接口访问。"
    if is_sensitive_path(full):
        return "错误：出于安全原因，禁止读取敏感文件。"

    if not full.exists():
        cwd_hint = f"（当前目录: {_working_directory}）" if _working_directory != "." else ""
        return f"错误：文件不存在 {path}{cwd_hint}（已搜索 input/, src/, output/ 等目录）"

    # Record file open in workspace context
    ws = _get_workspace_service()
    if ws is not None:
        try:
            ws.record_open(str(full.relative_to(ROOT)))
        except ValueError:
            pass

    suffix = full.suffix.lower()

    # Word document (.docx)
    if suffix == ".docx":
        try:
            from docx import Document
            doc = Document(str(full))
            paragraphs = [p.text for p in doc.paragraphs]
            tables_text = []
            for table in doc.tables:
                for row in table.rows:
                    cells = [cell.text for cell in row.cells]
                    tables_text.append(" | ".join(cells))
            all_text = "\n".join(paragraphs)
            if tables_text:
                all_text += "\n\n[表格内容]\n" + "\n".join(tables_text)
            return redact_sensitive_text(all_text.strip()) if all_text.strip() else "文档内容为空"
        except ImportError:
            return "错误：python-docx 未安装，无法读取 .docx 文件。请运行: pip install python-docx"
        except Exception as e:
            return f"错误：读取 .docx 文件失败: {str(e)}"

    # PowerPoint (.pptx)
    if suffix == ".pptx":
        try:
            from pptx import Presentation
            prs = Presentation(str(full))
            slides_text = []
            for i, slide in enumerate(prs.slides, 1):
                slide_lines = [f"--- 幻灯片 {i} ---"]
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for paragraph in shape.text_frame.paragraphs:
                            text = paragraph.text.strip()
                            if text:
                                slide_lines.append(text)
                    if shape.has_table:
                        table = shape.table
                        for row in table.rows:
                            cells = [cell.text for cell in row.cells]
                            slide_lines.append(" | ".join(cells))
                slides_text.append("\n\n".join(slide_lines))
            return redact_sensitive_text("\n\n".join(slides_text)) if slides_text else "PPT 内容为空"
        except ImportError:
            return "错误：python-pptx 未安装，无法读取 .pptx 文件。请运行: pip install python-pptx"
        except Exception as e:
            return f"错误：读取 .pptx 文件失败: {str(e)}"

    # Plain text
    try:
        return redact_sensitive_text(full.read_text(encoding="utf-8"))
    except UnicodeDecodeError:
        try:
            return redact_sensitive_text(full.read_text(encoding="gbk"))
        except Exception:
            return f"错误：无法解码文件 {path}，请确保文件是文本格式"


def write_file(
    path: str,
    content: str,
    mode: str = "overwrite",
    exact: bool = False,
) -> str:
    """Write content to a file.

    Args:
        path: file path (relative to current directory or project root)
        content: content to write
        mode: "overwrite" (default) or "append"

    Returns:
        Operation result description
    """
    if is_internal_storage_path(path):
        return "错误：PROTECTED_INTERNAL_PATH：Agent 内部存储只能通过专用接口访问。"
    if is_sensitive_path(path):
        return "错误：出于安全原因，禁止写入敏感文件。"

    if Path(str(path)).suffix.lower() in {".docx", ".xlsx", ".xls", ".pptx"}:
        return (
            "错误：Office 二进制文件不能通过文本 write_file 生成；"
            "请在可用沙箱中运行生成脚本，或先生成可审阅的 Python 脚本。"
        )

    try:
        # A write target is an explicit destination.  Fuzzy discovery is safe
        # for reads, but can redirect a new file to an unrelated existing
        # file (for example output/calc.py -> output/solution.py).
        requested = Path(str(path))
        candidate = requested if requested.is_absolute() else ROOT / requested
        full = _ensure_workspace_path(
            candidate if exact else _resolve_path(path),
            str(path),
        )
    except PermissionError as e:
        return f"错误：{e}"
    if is_internal_storage_path(full):
        return "错误：PROTECTED_INTERNAL_PATH：Agent 内部存储只能通过专用接口访问。"
    if is_sensitive_path(full):
        return "错误：出于安全原因，禁止写入敏感文件。"
    full.parent.mkdir(parents=True, exist_ok=True)

    if mode == "append":
        with open(full, "a", encoding="utf-8") as f:
            f.write(content)
        result = f"已追加内容到 {path}"
    else:
        full.write_text(content, encoding="utf-8")
        result = f"已写入 {path}"

    # Record file edit in workspace context
    ws = _get_workspace_service()
    if ws is not None:
        try:
            ws.record_edit(str(full.relative_to(ROOT)))
        except ValueError:
            pass

    return result


def copy_file(
    source: str,
    destination: str,
    exact: bool = True,
) -> str:
    """Copy one workspace file to another workspace path.

    This is the canonical ``filesystem.copy`` primitive.  It never falls back
    to shell execution and never accepts paths outside the active workspace.
    """
    for value in (source, destination):
        error = _protect_operation_path(value, label="复制")
        if error:
            return error
    try:
        source_path = _resolve_operation_path(source, exact=exact)
        destination_path = _resolve_operation_path(destination, exact=True)
    except PermissionError as exc:
        return f"错误：{exc}"
    if source_path == destination_path:
        return "错误：FILE_OPERATION_FAILED：复制源和目标不能相同。"
    if not source_path.exists() or not source_path.is_file():
        return f"错误：FILE_OPERATION_FAILED：复制源文件不存在: {source}"
    if is_internal_storage_path(source_path) or is_sensitive_path(source_path):
        return "错误：PROTECTED_INTERNAL_PATH：禁止复制内部或敏感文件。"
    if is_internal_storage_path(destination_path) or is_sensitive_path(destination_path):
        return "错误：PROTECTED_INTERNAL_PATH：禁止写入内部或敏感文件。"
    try:
        destination_path.parent.mkdir(parents=True, exist_ok=True)
        shutil.copy2(source_path, destination_path)
        _record_edit(destination_path)
    except OSError as exc:
        return f"错误：FILE_OPERATION_FAILED：复制失败: {exc}"
    return f"已复制 {source} 到 {destination}"


def move_file(
    source: str,
    destination: str,
    exact: bool = True,
) -> str:
    """Move one workspace file to another workspace path."""
    for value in (source, destination):
        error = _protect_operation_path(value, label="移动")
        if error:
            return error
    try:
        source_path = _resolve_operation_path(source, exact=exact)
        destination_path = _resolve_operation_path(destination, exact=True)
    except PermissionError as exc:
        return f"错误：{exc}"
    if source_path == destination_path:
        return "错误：FILE_OPERATION_FAILED：移动源和目标不能相同。"
    if not source_path.exists() or not source_path.is_file():
        return f"错误：FILE_OPERATION_FAILED：移动源文件不存在: {source}"
    if is_internal_storage_path(source_path) or is_sensitive_path(source_path):
        return "错误：PROTECTED_INTERNAL_PATH：禁止移动内部或敏感文件。"
    if is_internal_storage_path(destination_path) or is_sensitive_path(destination_path):
        return "错误：PROTECTED_INTERNAL_PATH：禁止写入内部或敏感文件。"
    try:
        destination_path.parent.mkdir(parents=True, exist_ok=True)
        shutil.move(str(source_path), str(destination_path))
        _record_edit(destination_path)
    except OSError as exc:
        return f"错误：FILE_OPERATION_FAILED：移动失败: {exc}"
    return f"已移动 {source} 到 {destination}"


def delete_file(path: str, exact: bool = True) -> str:
    """Delete one regular workspace file and verify the target is absent."""
    error = _protect_operation_path(path, label="删除")
    if error:
        return error
    try:
        full = _resolve_operation_path(path, exact=exact)
    except PermissionError as exc:
        return f"错误：{exc}"
    if is_internal_storage_path(full) or is_sensitive_path(full):
        return "错误：PROTECTED_INTERNAL_PATH：禁止删除内部或敏感文件。"
    if not full.exists():
        return f"错误：FILE_OPERATION_FAILED：文件不存在: {path}"
    if not full.is_file():
        return f"错误：FILE_OPERATION_FAILED：只允许删除文件，不能删除目录: {path}"
    try:
        full.unlink()
        _record_edit(full)
    except OSError as exc:
        return f"错误：FILE_OPERATION_FAILED：删除失败: {exc}"
    return f"已删除 {path}"


def list_directory(path: str = ".") -> str:
    """List directory contents.

    Args:
        path: directory path (relative to current directory or project root)

    Returns:
        Directory listing
    """
    if _working_directory != "." and not path.startswith("/"):
        cwd_full = (ROOT / _working_directory / path).resolve()
        if cwd_full.is_dir():
            full = cwd_full
        else:
            full = _resolve_path(path)
    else:
        full = _resolve_path(path)

    try:
        full = _ensure_workspace_path(full, str(path))
    except PermissionError as e:
        return f"错误：{e}"

    if not full.exists() or not full.is_dir():
        cwd_hint = f"（当前目录: {_working_directory}）" if _working_directory != "." else ""
        return f"错误：{path} 不是有效目录{cwd_hint}"

    items = [
        f"{p.name}{'/' if p.is_dir() else ''}"
        for p in full.iterdir()
        if not is_sensitive_path(p)
    ]
    display_path = str(full.relative_to(ROOT)) if str(full).startswith(str(ROOT)) else str(full)
    result = f"📁 {display_path}/ ({len(items)} 项)\n"
    result += "\n".join(sorted(items))
    return result


def find_file(name: str) -> str:
    """Find files by name (recursive search). Delegates to Workspace.

    Args:
        name: file name or pattern

    Returns:
        Matching file paths
    """
    ws = _get_workspace_service()
    if ws is not None:
        matches = ws.find(name)
        if matches:
            matches = [m for m in matches if not is_sensitive_path(m.path)]
        if matches:
            result_str = "\n".join(
                f"  [{m.score:.2f}] {m.path}  ({m.reason})"
                for m in matches[:10]
            )
            if len(matches) > 10:
                result_str += f"\n  ... 及另外 {len(matches) - 10} 个匹配"
            return f"找到 {len(matches)} 个文件:\n{result_str}"

    # Filesystem fallback when Workspace is unavailable
    results = []
    for p in ROOT.rglob(name):
        if (
            p.is_file()
            and not is_sensitive_path(p)
            and not any(part.startswith(".") for part in p.parts)
        ):
            try:
                rel = p.relative_to(ROOT)
                results.append(str(rel))
            except ValueError:
                results.append(str(p))

    if results:
        result_str = "\n".join(results[:5])
        if len(results) > 5:
            result_str += f"\n... 及另外 {len(results) - 5} 个匹配"
        return f"找到 {len(results)} 个文件:\n{result_str}"

    return f"未找到文件: {name}"


def clear_path_cache() -> str:
    """Clear path resolution cache."""
    _path_cache.clear()
    return "路径缓存已清除"


# ── Registration ──

registry.register(read_file, category="filesystem", tags=["filesystem", "read"])
registry.register(write_file, category="filesystem", tags=["filesystem", "write"])
registry.register(copy_file, category="filesystem", tags=["filesystem", "copy"])
registry.register(move_file, category="filesystem", tags=["filesystem", "move"])
registry.register(delete_file, category="filesystem", tags=["filesystem", "delete"])
registry.register(list_directory, category="filesystem", tags=["filesystem", "list"])
registry.register(set_working_directory, category="filesystem", tags=["filesystem", "navigate"])
registry.register(find_file, category="filesystem", tags=["filesystem", "search"])
registry.register(clear_path_cache, category="filesystem", tags=["filesystem", "debug"])

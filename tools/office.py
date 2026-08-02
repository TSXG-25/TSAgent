# tools/office.py
"""Office document creation tools — PPT and DOCX generation.

Provides:
- create_pptx(title, slides) — Beautiful PowerPoint presentation
- create_docx(title, sections) — Well-formatted Word document
"""
from pathlib import Path
from agent.registry.tool_registry import registry

ROOT = Path(__file__).parent.parent.resolve()
OUTPUT_DIR = ROOT / "output"
OUTPUT_DIR.mkdir(parents=True, exist_ok=True)

# ──────────────────────────────────────────────
# PPTX generation
# ──────────────────────────────────────────────


def _hex_to_rgb(hex_color: str):
    """Convert hex color string to pptx RGBColor."""
    from pptx.util import Pt, Inches, Emu
    from pptx.dml.color import RGBColor

    hex_color = hex_color.lstrip("#")
    return RGBColor(int(hex_color[:2], 16), int(hex_color[2:4], 16), int(hex_color[4:6], 16))


_COLOR_SCHEMES = {
    "modern": {
        "bg": "1E1E2E",           # dark navy bg
        "accent": "89B4FA",       # soft blue accent
        "title": "CDD6F4",        # light text
        "body": "A6ADC8",         # muted text
        "heading_bg": "313244",   # card background
    },
    "corporate": {
        "bg": "FFFFFF",
        "accent": "1A73E8",
        "title": "202124",
        "body": "5F6368",
        "heading_bg": "E8F0FE",
    },
    "gradient": {
        "bg": "0F0F23",
        "accent": "FF6B6B",
        "title": "FFFFFF",
        "body": "C0C0D0",
        "heading_bg": "1A1A3E",
    },
}


def create_pptx(title: str, slides: list, theme: str = "modern",
                output_path: str = "") -> str:
    """创建精美的 PowerPoint 演示文稿。

    支持标题、正文要点、引用、表格等多种幻灯片布局。
    自动应用专业配色主题。

    Args:
        title: 演示文稿标题（将作为第一页标题）
        slides: 幻灯片内容列表。每项为 dict，格式：
            {"type": "title|content|two_column|quote|table",
             "heading": "段落标题",
             "body": ["要点1", "要点2", ...],  # 或
             "left": [...], "right": [...],   # two_column
             "quote": "引用内容", "source": "来源",  # quote
             "headers": ["列1","列2"], "rows": [["a","b"],...],  # table
             "note": "备注文字（可选）"}
        theme: 配色主题 "modern" | "corporate" | "gradient"（默认 "modern"）
        output_path: 输出路径（相对于项目根目录，留空自动生成）

    Returns:
        操作结果描述（含文件路径）
    """
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt, Emu
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
        from pptx.enum.shapes import MSO_SHAPE
    except ImportError:
        return "错误：python-pptx 未安装。请运行: pip install python-pptx"

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    colors = _COLOR_SCHEMES.get(theme, _COLOR_SCHEMES["modern"])

    # ── Helper: add text box with styling ──
    def _add_textbox(slide, left, top, width, height, text, font_size=18,
                     bold=False, color=None, alignment=PP_ALIGN.LEFT,
                     font_name="Microsoft YaHei"):
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.bold = bold
        p.font.color.rgb = _hex_to_rgb(color or colors["body"])
        p.font.name = font_name
        p.alignment = alignment
        return tf

    def _add_bullet_list(tf, items: list, font_size=16, color=None):
        """Add bullet points to an existing text frame."""
        for item in items:
            p = tf.add_paragraph()
            p.text = item
            p.font.size = Pt(font_size)
            p.font.color.rgb = _hex_to_rgb(color or colors["body"])
            p.font.name = "Microsoft YaHei"
            p.space_after = Pt(6)
            # bullet character
            p.level = 0

    # ── Slide 1: Title slide ──
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout

    # Background rectangle
    bg_shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
        prs.slide_width, prs.slide_height
    )
    bg_shape.fill.solid()
    bg_shape.fill.fore_color.rgb = _hex_to_rgb(colors["bg"])
    bg_shape.line.fill.background()

    # Accent bar on left
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.3), prs.slide_height
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = _hex_to_rgb(colors["accent"])
    bar.line.fill.background()

    # Title
    _add_textbox(slide, Inches(1.5), Inches(2.0), Inches(10), Inches(1.5),
                 title, font_size=44, bold=True, color=colors["title"],
                 alignment=PP_ALIGN.LEFT)

    # Subtitle line
    from datetime import datetime
    subtitle = f"生成时间: {datetime.now().strftime('%Y-%m-%d %H:%M')}"
    _add_textbox(slide, Inches(1.5), Inches(3.8), Inches(8), Inches(0.8),
                 subtitle, font_size=20, color=colors["body"],
                 alignment=PP_ALIGN.LEFT)

    # ── Content slides ──
    for slide_data in slides:
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        # Background
        bg_shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
            prs.slide_width, prs.slide_height
        )
        bg_shape.fill.solid()
        bg_shape.fill.fore_color.rgb = _hex_to_rgb(colors["bg"])
        bg_shape.line.fill.background()

        # Top accent bar
        top_bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
            prs.slide_width, Inches(0.08)
        )
        top_bar.fill.solid()
        top_bar.fill.fore_color.rgb = _hex_to_rgb(colors["accent"])
        top_bar.line.fill.background()

        slide_type = slide_data.get("type", "content")
        heading = slide_data.get("heading", "")

        # Section number
        idx = slides.index(slide_data) + 1

        # Heading with number badge
        heading_badge = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.8), Inches(0.4), Inches(0.6), Inches(0.6)
        )
        heading_badge.fill.solid()
        heading_badge.fill.fore_color.rgb = _hex_to_rgb(colors["accent"])
        heading_badge.line.fill.background()
        tf = heading_badge.text_frame
        tf.paragraphs[0].text = str(idx).zfill(2)
        tf.paragraphs[0].font.size = Pt(22)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = _hex_to_rgb(colors["bg"])
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        tf.word_wrap = False

        if heading:
            _add_textbox(slide, Inches(1.6), Inches(0.35), Inches(10), Inches(0.8),
                         heading, font_size=32, bold=True, color=colors["title"])

        content_top = Inches(1.5)

        if slide_type == "content":
            body = slide_data.get("body", [])
            if body:
                # Content card background
                card = slide.shapes.add_shape(
                    MSO_SHAPE.ROUNDED_RECTANGLE,
                    Inches(0.8), content_top, Inches(11.5), Inches(5.0)
                )
                card.fill.solid()
                card.fill.fore_color.rgb = _hex_to_rgb(colors["heading_bg"])
                card.line.fill.background()

                tf = _add_textbox(slide, Inches(1.2), Inches(content_top + 0.3),
                                  Inches(10.5), Inches(0.5),
                                  "▎ 详细内容", font_size=18, bold=True,
                                  color=colors["accent"])
                tf2 = _add_textbox(slide, Inches(1.2), Inches(content_top + 1.0),
                                   Inches(10.5), Inches(3.5),
                                   "", font_size=16, color=colors["body"])
                _add_bullet_list(tf2, body, font_size=16)

        elif slide_type == "two_column":
            left = slide_data.get("left", [])
            right = slide_data.get("right", [])

            # Left column card
            card_l = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(0.8), content_top, Inches(5.5), Inches(5.0)
            )
            card_l.fill.solid()
            card_l.fill.fore_color.rgb = _hex_to_rgb(colors["heading_bg"])
            card_l.line.fill.background()

            heading_l = slide_data.get("heading_left", "左栏")
            _add_textbox(slide, Inches(1.1), Inches(content_top + 0.2),
                         Inches(5), Inches(0.5),
                         heading_l, font_size=18, bold=True, color=colors["accent"])
            tf_l = _add_textbox(slide, Inches(1.1), Inches(content_top + 0.8),
                                Inches(5), Inches(4),
                                "", font_size=15, color=colors["body"])
            _add_bullet_list(tf_l, left, font_size=15)

            # Right column card
            card_r = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(6.8), content_top, Inches(5.5), Inches(5.0)
            )
            card_r.fill.solid()
            card_r.fill.fore_color.rgb = _hex_to_rgb(colors["heading_bg"])
            card_r.line.fill.background()

            heading_r = slide_data.get("heading_right", "右栏")
            _add_textbox(slide, Inches(7.1), Inches(content_top + 0.2),
                         Inches(5), Inches(0.5),
                         heading_r, font_size=18, bold=True, color=colors["accent"])
            tf_r = _add_textbox(slide, Inches(7.1), Inches(content_top + 0.8),
                                Inches(5), Inches(4),
                                "", font_size=15, color=colors["body"])
            _add_bullet_list(tf_r, right, font_size=15)

        elif slide_type == "quote":
            quote = slide_data.get("quote", "")
            source = slide_data.get("source", "")

            # Quote card
            card = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(1.5), Inches(2.0), Inches(10), Inches(3.5)
            )
            card.fill.solid()
            card.fill.fore_color.rgb = _hex_to_rgb(colors["heading_bg"])
            card.line.fill.background()

            # Quote mark
            _add_textbox(slide, Inches(2.0), Inches(2.3), Inches(1), Inches(1),
                         "❝", font_size=48, color=colors["accent"])

            _add_textbox(slide, Inches(2.5), Inches(2.5), Inches(8), Inches(1.5),
                         quote, font_size=24, color=colors["title"])
            if source:
                _add_textbox(slide, Inches(2.5), Inches(4.2), Inches(8), Inches(0.5),
                             f"— {source}", font_size=16, color=colors["body"])

        elif slide_type == "table":
            headers = slide_data.get("headers", [])
            rows = slide_data.get("rows", [])

            if headers:
                table_shape = slide.shapes.add_table(
                    len(rows) + 1, len(headers),
                    Inches(0.8), content_top, Inches(11.5), Inches(0.4 * (len(rows) + 2))
                )
                table = table_shape.table

                # Header row
                for j, h in enumerate(headers):
                    cell = table.cell(0, j)
                    cell.text = h
                    for paragraph in cell.text_frame.paragraphs:
                        paragraph.font.size = Pt(14)
                        paragraph.font.bold = True
                        paragraph.font.color.rgb = _hex_to_rgb(colors["bg"])
                        paragraph.alignment = PP_ALIGN.CENTER
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = _hex_to_rgb(colors["accent"])

                # Data rows
                for i, row in enumerate(rows):
                    for j, val in enumerate(row):
                        cell = table.cell(i + 1, j)
                        cell.text = str(val)
                        for paragraph in cell.text_frame.paragraphs:
                            paragraph.font.size = Pt(12)
                            paragraph.font.color.rgb = _hex_to_rgb(colors["body"])
                            paragraph.alignment = PP_ALIGN.CENTER
                        if i % 2 == 0:
                            cell.fill.solid()
                            cell.fill.fore_color.rgb = _hex_to_rgb(colors["heading_bg"])

        # Speaker notes
        note = slide_data.get("note", "")
        if note:
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = note

    # ── Save ──
    if not output_path:
        safe_title = "".join(c if c.isalnum() or c in " _-" else "_" for c in title)
        output_path = f"output/{safe_title}_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx"

    full_path = (ROOT / output_path).resolve()
    full_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(full_path))

    return f"✅ PPT 已生成: {output_path}\n幻灯片数量: {len(slides) + 1} 页\n配色主题: {theme}"


# ──────────────────────────────────────────────
# DOCX generation
# ──────────────────────────────────────────────


def create_docx(title: str, sections: list, output_path: str = "") -> str:
    """创建格式精美的 Word 文档。

    支持标题、段落正文、要点列表、表格、引用等多种内容块。
    自动应用专业的排版和样式。

    Args:
        title: 文档标题
        sections: 文档内容块列表。每项为 dict，格式：
            {"type": "heading|paragraph|bullet|table|quote",
             "text": "内容文字",
             "level": 1-3（仅 heading 类型，默认 1）,
             "items": ["点1","点2"]（仅 bullet 类型）,
             "headers": ["列1","列2"], "rows": [["a","b"]]（仅 table 类型）,
             "source": "来源"（仅 quote 类型）}
        output_path: 输出路径（相对于项目根目录，留空自动生成）

    Returns:
        操作结果描述（含文件路径）
    """
    try:
        from docx import Document
        from docx.shared import Pt, Inches, Cm, RGBColor
        from docx.enum.text import WD_ALIGN_PARAGRAPH
        from docx.oxml.ns import qn
        from docx.oxml import OxmlElement
    except ImportError:
        return "错误：python-docx 未安装。请运行: pip install python-docx"

    doc = Document()

    # ── Style setup ──
    style = doc.styles["Normal"]
    font = style.font
    font.name = "Microsoft YaHei"
    font.size = Pt(11)
    font.color.rgb = RGBColor(0x33, 0x33, 0x33)
    style.paragraph_format.space_after = Pt(6)

    # Set CJK font
    rPr = style.element.get_or_add_rPr()
    rFonts = rPr.find(qn('w:rFonts'))
    if rFonts is None:
        rFonts = OxmlElement('w:rFonts')
        rPr.append(rFonts)
    rFonts.set(qn('w:eastAsia'), 'Microsoft YaHei')

    # Heading styles
    for level in range(1, 4):
        hs = doc.styles[f"Heading {level}"]
        hs.font.name = "Microsoft YaHei"
        hs.font.color.rgb = RGBColor(0x1A, 0x73, 0xE8)
        hs.font.bold = True
        if level == 1:
            hs.font.size = Pt(24)
            hs.paragraph_format.space_before = Pt(24)
            hs.paragraph_format.space_after = Pt(12)
        elif level == 2:
            hs.font.size = Pt(18)
            hs.paragraph_format.space_before = Pt(18)
            hs.paragraph_format.space_after = Pt(8)
        else:
            hs.font.size = Pt(14)
            hs.paragraph_format.space_before = Pt(12)
            hs.paragraph_format.space_after = Pt(6)

    # ── Title page ──
    # Add some spacing before title
    for _ in range(4):
        doc.add_paragraph("")

    title_para = doc.add_paragraph()
    title_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = title_para.add_run(title)
    run.font.size = Pt(32)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x1A, 0x73, 0xE8)
    run.font.name = "Microsoft YaHei"

    # Subtitle with date
    from datetime import datetime
    sub_para = doc.add_paragraph()
    sub_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = sub_para.add_run(f"生成时间: {datetime.now().strftime('%Y-%m-%d %H:%M')}")
    run.font.size = Pt(14)
    run.font.color.rgb = RGBColor(0x66, 0x66, 0x66)
    run.font.name = "Microsoft YaHei"

    # Divider line
    doc.add_paragraph("")
    divider = doc.add_paragraph()
    divider.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = divider.add_run("─" * 50)
    run.font.color.rgb = RGBColor(0x1A, 0x73, 0xE8)
    run.font.size = Pt(10)

    doc.add_page_break()

    # ── Content sections ──
    for section in sections:
        stype = section.get("type", "paragraph")

        if stype == "heading":
            level = section.get("level", 1)
            doc.add_heading(section.get("text", ""), level=level)

        elif stype == "paragraph":
            p = doc.add_paragraph(section.get("text", ""))
            # First line indent for readability
            p.paragraph_format.first_line_indent = Cm(0.75)

        elif stype == "bullet":
            items = section.get("items", [])
            for item in items:
                p = doc.add_paragraph(item, style="List Bullet")
                for run in p.runs:
                    run.font.name = "Microsoft YaHei"
                    run.font.size = Pt(11)

        elif stype == "quote":
            quote_text = section.get("text", "")
            source = section.get("source", "")

            # Indented quote
            p = doc.add_paragraph()
            p.paragraph_format.left_indent = Cm(1.5)
            p.paragraph_format.right_indent = Cm(1.5)
            run = p.add_run(f"❝ {quote_text}")
            run.font.italic = True
            run.font.size = Pt(12)
            run.font.color.rgb = RGBColor(0x55, 0x55, 0x55)

            if source:
                p2 = doc.add_paragraph()
                p2.paragraph_format.left_indent = Cm(1.5)
                run = p2.add_run(f"— {source}")
                run.font.size = Pt(10)
                run.font.color.rgb = RGBColor(0x88, 0x88, 0x88)

        elif stype == "table":
            headers = section.get("headers", [])
            rows = section.get("rows", [])

            if headers and rows:
                table = doc.add_table(rows=len(rows) + 1, cols=len(headers))
                table.style = "Light Grid Accent 1"

                # Header row
                for j, h in enumerate(headers):
                    cell = table.cell(0, j)
                    cell.text = h
                    for paragraph in cell.paragraphs:
                        for run in paragraph.runs:
                            run.font.bold = True
                            run.font.size = Pt(10)

                # Data rows
                for i, row in enumerate(rows):
                    for j, val in enumerate(row):
                        cell = table.cell(i + 1, j)
                        cell.text = str(val)
                        for paragraph in cell.paragraphs:
                            for run in paragraph.runs:
                                run.font.size = Pt(10)

        # Add spacing between sections
        doc.add_paragraph("")

    # ── Save ──
    if not output_path:
        safe_title = "".join(c if c.isalnum() or c in " _-" else "_" for c in title)
        output_path = f"output/{safe_title}_{datetime.now().strftime('%Y%m%d_%H%M%S')}.docx"

    full_path = (ROOT / output_path).resolve()
    full_path.parent.mkdir(parents=True, exist_ok=True)
    doc.save(str(full_path))

    return f"✅ Word 文档已生成: {output_path}\n内容块数量: {len(sections) * 2 + 1} 段"


# ── Register tools ──

registry.register(create_pptx, category="office", tags=["office", "ppt", "presentation"])
registry.register(create_docx, category="office", tags=["office", "docx", "word"])